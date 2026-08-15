from __future__ import annotations

import uuid
from dataclasses import dataclass
from pathlib import Path
from typing import Any


@dataclass(frozen=True)
class DocumentExtraction:
    text: str
    chunks: list[dict[str, Any]]


def _split_text(text: str, max_chars: int = 1600, overlap: int = 180) -> list[str]:
    normalized = "\n".join(line.strip() for line in text.splitlines() if line.strip())
    if not normalized:
        return []
    pieces: list[str] = []
    start = 0
    while start < len(normalized):
        end = min(len(normalized), start + max_chars)
        if end < len(normalized):
            boundary = max(normalized.rfind("\n", start, end), normalized.rfind("。", start, end))
            if boundary > start + max_chars // 2:
                end = boundary + 1
        pieces.append(normalized[start:end])
        if end >= len(normalized):
            break
        start = max(start + 1, end - overlap)
    return pieces


def _chunk(
    text: str,
    *,
    position: int,
    locator_type: str,
    page: int | None = None,
    slide: int | None = None,
    content_kind: str = "text",
    visual_path: str | None = None,
) -> dict[str, Any]:
    return {
        "id": str(uuid.uuid4()),
        "position": position,
        "text": text,
        "locator_type": locator_type,
        "page": page,
        "slide": slide,
        "content_kind": content_kind,
        "visual_path": visual_path,
        "metadata": {},
    }


def extract_document(path: Path, mime_type: str | None, output_dir: Path) -> DocumentExtraction:
    suffix = path.suffix.lower()
    output_dir.mkdir(parents=True, exist_ok=True)
    if suffix == ".pdf" or mime_type == "application/pdf":
        return _extract_pdf(path, output_dir)
    if suffix in {".pptx", ".ppt"}:
        return _extract_presentation(path, output_dir)
    if suffix in {".txt", ".md", ".csv"} or (mime_type or "").startswith("text/"):
        text = path.read_text(encoding="utf-8", errors="replace")
        chunks = [
            _chunk(piece, position=index, locator_type="chunk")
            for index, piece in enumerate(_split_text(text), start=1)
        ]
        return DocumentExtraction(text=text, chunks=chunks)
    return DocumentExtraction(text="", chunks=[])


def _extract_pdf(path: Path, output_dir: Path) -> DocumentExtraction:
    import fitz

    document = fitz.open(path)
    chunks: list[dict[str, Any]] = []
    full_text: list[str] = []
    position = 1
    try:
        for page_number, page in enumerate(document, start=1):
            text = page.get_text("text")
            full_text.append(f"[PDF page {page_number}]\n{text}")
            image_path = output_dir / f"page-{page_number}.png"
            page.get_pixmap(matrix=fitz.Matrix(1, 1), alpha=False).save(image_path)
            pieces = _split_text(text) or ["[This page contains visual content without extractable text]"]
            for piece in pieces:
                chunks.append(
                    _chunk(
                        piece,
                        position=position,
                        locator_type="page",
                        page=page_number,
                        content_kind="mixed" if not text.strip() else "text",
                        visual_path=str(image_path),
                    )
                )
                position += 1
    finally:
        document.close()
    return DocumentExtraction(text="\n\n".join(full_text), chunks=chunks)


def _extract_presentation(path: Path, output_dir: Path) -> DocumentExtraction:
    if path.suffix.lower() == ".ppt":
        return DocumentExtraction(text="", chunks=[])
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    presentation = Presentation(str(path))
    chunks: list[dict[str, Any]] = []
    full_text: list[str] = []
    position = 1
    for slide_number, slide in enumerate(presentation.slides, start=1):
        texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
        visual_paths: list[str] = []
        for image_number, shape in enumerate(
            (item for item in slide.shapes if item.shape_type == MSO_SHAPE_TYPE.PICTURE), start=1
        ):
            extension = shape.image.ext or "bin"
            image_path = output_dir / f"slide-{slide_number}-image-{image_number}.{extension}"
            image_path.write_bytes(shape.image.blob)
            visual_paths.append(str(image_path))
        text = "\n".join(texts)
        full_text.append(f"[PPT slide {slide_number}]\n{text}")
        pieces = _split_text(text) or ["[This slide contains visual content without extractable text]"]
        for piece in pieces:
            chunks.append(
                _chunk(
                    piece,
                    position=position,
                    locator_type="slide",
                    slide=slide_number,
                    content_kind="mixed" if visual_paths else "text",
                    visual_path=visual_paths[0] if visual_paths else None,
                )
            )
            chunks[-1]["metadata"] = {"visual_paths": visual_paths}
            position += 1
    return DocumentExtraction(text="\n\n".join(full_text), chunks=chunks)


def extract_text(path: Path, mime_type: str | None) -> str:
    """Backward-compatible text-only helper for integrations."""
    return extract_document(path, mime_type, path.parent / f"{path.name}.assets").text
