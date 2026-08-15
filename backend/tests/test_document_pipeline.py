from pathlib import Path

import fitz
import pytest
from pptx import Presentation

from app.database import Database
from app.documents import extract_document
from app.models import CourseCreate, SessionCreate
from app.providers.hash_embedding import HashEmbeddingProvider
from app.retrieval import RetrievalPolicy, SessionRetriever


def test_pdf_is_split_by_real_pages_and_keeps_visuals(tmp_path: Path):
    pdf_path = tmp_path / "book.pdf"
    document = fitz.open()
    first = document.new_page()
    first.insert_text((72, 72), "Limits and continuity are prerequisites.")
    second = document.new_page()
    second.insert_text((72, 72), "The mean value theorem needs differentiability.")
    document.save(pdf_path)
    document.close()

    extraction = extract_document(pdf_path, "application/pdf", tmp_path / "assets")

    assert {chunk["page"] for chunk in extraction.chunks} == {1, 2}
    assert all(chunk["locator_type"] == "page" for chunk in extraction.chunks)
    assert all(Path(chunk["visual_path"]).exists() for chunk in extraction.chunks)
    assert "[PDF page 2]" in extraction.text


def test_pptx_is_split_by_slide_and_extracts_text(tmp_path: Path):
    pptx_path = tmp_path / "lecture.pptx"
    presentation = Presentation()
    first = presentation.slides.add_slide(presentation.slide_layouts[1])
    first.shapes.title.text = "Rolle's theorem"
    first.placeholders[1].text = "Continuity on [a,b]"
    second = presentation.slides.add_slide(presentation.slide_layouts[1])
    second.shapes.title.text = "Lagrange theorem"
    second.placeholders[1].text = "Differentiability on (a,b)"
    presentation.save(pptx_path)

    extraction = extract_document(pptx_path, "application/vnd.openxmlformats-officedocument.presentationml.presentation", tmp_path / "assets")

    assert [chunk["slide"] for chunk in extraction.chunks] == [1, 2]
    assert "[PPT slide 2]" in extraction.text
    assert "Differentiability" in extraction.chunks[1]["text"]


@pytest.mark.anyio
async def test_session_retrieval_uses_different_reconstruction_and_learning_policies(tmp_path: Path):
    database = Database(tmp_path / "retrieval.sqlite3")
    course = database.create_course(CourseCreate(name="Calculus"))
    session = database.create_session(course["id"], SessionCreate(title="Mean value theorem"))
    slides = database.add_resource(
        session["id"], type="slides", evidence_level="official", name="lecture.pptx", extracted_text="theorem"
    )
    textbook = database.add_resource(
        session["id"], type="textbook", evidence_level="official", name="book.pdf", extracted_text="theorem"
    )
    for position, resource in enumerate((slides, textbook), start=1):
        database.replace_document_chunks(
            resource["id"],
            [
                {
                    "id": f"chunk-{position}",
                    "position": 1,
                    "text": "mean value theorem conditions",
                    "locator_type": "slide" if resource["type"] == "slides" else "page",
                    "slide": 1 if resource["type"] == "slides" else None,
                    "page": 1 if resource["type"] == "textbook" else None,
                }
            ],
        )

    retriever = SessionRetriever(database, HashEmbeddingProvider())
    await retriever.index_resource(slides["id"])
    await retriever.index_resource(textbook["id"])

    reconstruction = await retriever.retrieve(session["id"], "mean value theorem", RetrievalPolicy.RECONSTRUCTION)
    learning = await retriever.retrieve(session["id"], "mean value theorem", RetrievalPolicy.LEARNING)

    assert reconstruction[0]["resource_type"] == "slides"
    assert learning[0]["resource_type"] == "textbook"
