from __future__ import annotations

import wave
from pathlib import Path

import fitz
from fastapi.testclient import TestClient
from pptx import Presentation

from app.config import Settings
from app.database import Database
from app.main import create_app
from app.models import (
    Confidence,
    EvaluationResult,
    KnowledgePointDraft,
    KnowledgePointEvaluation,
    LearningStepDraft,
    MasteryEvidenceType,
    QuestionDraft,
    ReconstructionDraft,
    RemediationDraft,
    SourceRef,
    TimelineItem,
    TranscriptSegment,
)


class GroundedE2EProvider:
    """Deterministic evaluator that still requires real parsed and timestamped evidence."""

    requires_external_upload = False

    @staticmethod
    def _resource(evidence: list[dict], kind: str) -> dict:
        return next(item for item in evidence if item["type"] == kind)

    async def transcribe(self, path: str, mime_type: str | None) -> list[TranscriptSegment]:
        assert mime_type == "audio/wav"
        with wave.open(path, "rb") as recording:
            assert recording.getnchannels() == 1
            assert recording.getnframes() / recording.getframerate() == 5400
        return [
            TranscriptSegment(start_time=300, end_time=330, text="Raft elects one leader per term."),
            TranscriptSegment(start_time=1800, end_time=1830, text="A majority replicates each log entry."),
            TranscriptSegment(start_time=4200, end_time=4230, text="Committed entries survive leader changes."),
        ]

    async def analyze_session(self, session: dict, evidence: list[dict]) -> ReconstructionDraft:
        audio = self._resource(evidence, "audio")
        slides = self._resource(evidence, "slides")
        textbook = self._resource(evidence, "textbook")
        assert len(slides["chunks"]) == 40
        assert len(textbook["chunks"]) == 80
        transcript = SourceRef(
            resource_id=audio["id"],
            label=audio["name"],
            locator="30:00–30:30",
            locator_type="transcript",
            start_time=1800,
            end_time=1830,
        )
        slide = SourceRef(
            resource_id=slides["id"],
            label=slides["name"],
            locator="slide 20",
            locator_type="slide",
            slide=20,
        )
        page = SourceRef(
            resource_id=textbook["id"],
            label=textbook["name"],
            locator="page 40",
            locator_type="page",
            page=40,
        )
        election = KnowledgePointDraft(
            title="Raft leader election",
            description="Terms, votes, randomized timeouts, and majority election.",
            importance=5,
            expected_mastery_level=2,
            confidence=Confidence.CONFIRMED,
            sources=[slide, transcript],
        )
        replication = KnowledgePointDraft(
            title="Raft log replication",
            description="Majority replication, commit index, and leader completeness.",
            prerequisites=[election.title],
            importance=5,
            expected_mastery_level=2,
            confidence=Confidence.CONFIRMED,
            sources=[transcript, page],
        )
        return ReconstructionDraft(
            title=session["title"],
            summary="The lecture connected leader election to safe replicated-log commitment.",
            topics=[election.title, replication.title],
            timeline=[
                TimelineItem(
                    start_time=1800,
                    end_time=1830,
                    title="Majority replication",
                    summary="The teacher derived the commit rule.",
                    confidence=Confidence.CONFIRMED,
                    sources=[transcript],
                )
            ],
            teacher_emphasis=["A single successful answer is not sufficient evidence of mastery."],
            confirmed=[election.title, replication.title],
            knowledge_points=[election, replication],
            learning_path=[
                LearningStepDraft(
                    position=1,
                    title="Build the election state machine",
                    brief_explanation="Track terms and votes.",
                    full_explanation="A candidate increments its term and needs votes from a majority.",
                    knowledge_point_titles=[election.title],
                    estimated_minutes=12,
                    confidence=Confidence.CONFIRMED,
                    sources=[slide],
                ),
                LearningStepDraft(
                    position=2,
                    title="Prove log commitment safety",
                    brief_explanation="Connect quorum overlap to leader completeness.",
                    full_explanation="A committed entry is present in the quorum that elects a later leader.",
                    knowledge_point_titles=[replication.title],
                    estimated_minutes=18,
                    confidence=Confidence.CONFIRMED,
                    sources=[page, transcript],
                ),
            ],
        )

    async def generate_questions(
        self, session: dict, evidence: list[dict], knowledge_points: list[dict]
    ) -> list[QuestionDraft]:
        del session
        slides = self._resource(evidence, "slides")
        source = SourceRef(
            resource_id=slides["id"],
            label=slides["name"],
            locator="slide 20",
            locator_type="slide",
            slide=20,
        )
        questions: list[QuestionDraft] = []
        for point in knowledge_points:
            for level in ("understanding", "application"):
                questions.append(
                    QuestionDraft(
                        knowledge_point_titles=[point["title"]],
                        prompt=f"{point['title']} — demonstrate {level} with a quorum example.",
                        level=level,
                        question_type="application" if level == "application" else "diagnostic",
                        expected_mastery_level=point["expected_mastery"],
                        reference_answer="Explain the majority rule and its safety consequence.",
                        rubric=["uses a majority", "connects the rule to safety"],
                        source_refs=[source],
                    )
                )
        return questions

    async def evaluate_answer(self, question: dict, answer: str, evidence: list[dict]) -> EvaluationResult:
        del evidence
        assert "majority" in answer.lower()
        title = question["prompt"].split(" — ", 1)[0]
        evidence_type = (
            MasteryEvidenceType.APPLICATION
            if question["level"] == "application"
            else MasteryEvidenceType.UNDERSTANDING
        )
        return EvaluationResult(
            score=0.95,
            verdict="mastered",
            met_criteria=["uses a majority", "connects the rule to safety"],
            feedback="The answer is grounded in the quorum argument.",
            point_results=[
                KnowledgePointEvaluation(
                    knowledge_point_title=title,
                    score=0.95,
                    evidence_type=evidence_type,
                    feedback="Strong evidence.",
                )
            ],
        )

    async def remediate(self, knowledge_point: dict, reason: str, evidence: list[dict]) -> RemediationDraft:
        del reason
        textbook = self._resource(evidence, "textbook")
        return RemediationDraft(
            knowledge_point_title=knowledge_point["title"],
            diagnosis="The quorum-overlap link was missing.",
            simpler_explanation="Any two majorities overlap in at least one server.",
            analogy="Two majority committees must share at least one member.",
            worked_example="In five servers, any two groups of three overlap.",
            quick_check="Why can two disjoint groups of three not exist among five servers?",
            sources=[
                SourceRef(
                    resource_id=textbook["id"],
                    label=textbook["name"],
                    locator="page 40",
                    locator_type="page",
                    page=40,
                )
            ],
        )


def _write_long_recording(path: Path) -> None:
    sample_rate = 1000
    duration_seconds = 5400
    with wave.open(str(path), "wb") as recording:
        recording.setnchannels(1)
        recording.setsampwidth(1)
        recording.setframerate(sample_rate)
        block = bytes([128]) * sample_rate * 60
        for _ in range(duration_seconds // 60):
            recording.writeframesraw(block)


def _write_slides(path: Path) -> None:
    presentation = Presentation()
    for number in range(1, 41):
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = f"Raft consensus · slide {number}"
        slide.placeholders[1].text = (
            "Leader election, majority votes, log replication, commit index, and safety. "
            f"Worked scenario {number}."
        )
    presentation.save(path)


def _write_textbook(path: Path) -> None:
    document = fitz.open()
    for number in range(1, 81):
        page = document.new_page()
        page.insert_text(
            (72, 72),
            f"Raft textbook page {number}\nLeader election and majority log replication preserve safety.",
            fontsize=12,
        )
    document.save(path)
    document.close()


def _upload(client: TestClient, session_id: str, path: Path, mime_type: str, data: dict) -> dict:
    with path.open("rb") as handle:
        response = client.post(
            f"/sessions/{session_id}/resources/upload",
            data=data,
            files={"file": (path.name, handle, mime_type)},
        )
    assert response.status_code == 201, response.text
    return response.json()


def test_realistic_long_session_runs_from_capture_to_mastery(tmp_path: Path):
    audio_path = tmp_path / "lecture-90-minutes.wav"
    slides_path = tmp_path / "lecture-40-slides.pptx"
    textbook_path = tmp_path / "distributed-systems-80-pages.pdf"
    _write_long_recording(audio_path)
    _write_slides(slides_path)
    _write_textbook(textbook_path)

    settings = Settings(tmp_path / "runtime", "grounded", "grounded", None, "local", "local", "local")
    provider = GroundedE2EProvider()
    app = create_app(settings, Database(tmp_path / "e2e.sqlite3"), provider, provider)
    client = TestClient(app)

    course = client.post(
        "/courses",
        json={"name": "Distributed Systems", "semester": "2026 Fall"},
    ).json()
    session = client.post(
        f"/courses/{course['id']}/sessions",
        json={
            "title": "Raft consensus: election, replication, and safety",
            "starts_at": "2026-08-15T08:00:00+08:00",
            "ends_at": "2026-08-15T09:40:00+08:00",
            "notes": "The student missed this lecture and needs a from-zero recovery path.",
        },
    ).json()

    audio = _upload(
        client,
        session["id"],
        audio_path,
        "audio/wav",
        {
            "resource_type": "audio",
            "evidence_level": "classroom",
            "coverage": "0.9",
            "quality": "0.95",
            "duration_seconds": "5400",
            "start_offset": "0",
            "end_offset": "5400",
            "session_duration": "6000",
        },
    )
    transcript = client.post(f"/resources/{audio['id']}/transcribe", json={})
    assert transcript.status_code == 200
    assert transcript.json()["segments"][1]["global_start"] == 1800

    slides = _upload(
        client,
        session["id"],
        slides_path,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        {"resource_type": "slides", "evidence_level": "official"},
    )
    textbook = _upload(
        client,
        session["id"],
        textbook_path,
        "application/pdf",
        {"resource_type": "textbook", "evidence_level": "official"},
    )
    assert len(slides["chunks"]) == 40
    assert len(textbook["chunks"]) == 80

    manifest = client.get(
        f"/sessions/{session['id']}/consent-manifest", params={"operation": "analysis"}
    ).json()
    assert manifest["external"] is False
    assert {item["id"] for item in manifest["resources"]} == {audio["id"], slides["id"], textbook["id"]}

    reconstruction_results = client.post(
        f"/sessions/{session['id']}/retrieve",
        json={"query": "Raft leader election majority replication", "policy": "reconstruction"},
    ).json()
    learning_results = client.post(
        f"/sessions/{session['id']}/retrieve",
        json={"query": "Raft leader election majority replication", "policy": "learning"},
    ).json()
    assert any(item["resource_id"] == slides["id"] for item in reconstruction_results)
    assert any(item["resource_id"] == textbook["id"] for item in learning_results)

    analysis_job = client.post(
        f"/sessions/{session['id']}/jobs", json={"kind": "analysis"}
    ).json()
    assert client.get(f"/jobs/{analysis_job['id']}").json()["status"] == "succeeded"
    detail = client.get(f"/sessions/{session['id']}").json()
    assert detail["reconstruction_score"] >= 80
    assert detail["learning_coverage"] >= 90
    assert len(detail["knowledge_points"]) == 2
    assert detail["debts"][0]["status"] == "unseen"

    remediation = client.post(
        f"/knowledge-points/{detail['knowledge_points'][1]['id']}/remediation",
        json={"reason": "I cannot connect quorum overlap to safety."},
    )
    assert remediation.status_code == 201
    assert remediation.json()["payload"]["sources"][0]["page"] == 40

    assessment_job = client.post(
        f"/sessions/{session['id']}/jobs", json={"kind": "assessment"}
    ).json()
    assert client.get(f"/jobs/{assessment_job['id']}").json()["status"] == "succeeded"
    questions = client.get(f"/sessions/{session['id']}/assessment").json()
    assert len(questions) == 4

    for question in questions:
        result = client.post(
            f"/questions/{question['id']}/answer",
            json={"answer": "A majority quorum overlaps later majorities, which preserves safety."},
        )
        assert result.status_code == 200, result.text

    for step in detail["learning_steps"]:
        assert client.post(f"/learning-steps/{step['id']}/complete").json()["completed"] is True

    completed = client.get(f"/sessions/{session['id']}").json()
    assert completed["status"] == "complete"
    assert all(debt["status"] == "mastered" for debt in completed["debts"])
    assert client.get("/home").json()["open_debt_count"] == 0
